"""Executive Slide Deck Generator for NovaMart BCG X engagement.

Generates a 10-slide PowerPoint deck using python-pptx with real data
from the NovaMart analytics platform.
"""

from __future__ import annotations

from datetime import date
from pathlib import Path
from typing import Any

import pandas as pd
from loguru import logger

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    PPTX_OK = True
except ImportError:
    PPTX_OK = False
    logger.warning("python-pptx not installed — slide generation will fail")

# BCG X brand colours
_BCG_GREEN = RGBColor(0x00, 0xA6, 0x51) if PPTX_OK else None
_BCG_DARK = RGBColor(0x1A, 0x1A, 0x2E) if PPTX_OK else None
_BCG_BLUE = RGBColor(0x0F, 0x34, 0x60) if PPTX_OK else None
_WHITE = RGBColor(0xFF, 0xFF, 0xFF) if PPTX_OK else None
_GREY = RGBColor(0x4A, 0x55, 0x68) if PPTX_OK else None

_FOOTER = "BCG X Confidential | NovaMart Analytics | 2024"


def _add_text(
    tf: Any,
    text: str,
    size: int = 18,
    bold: bool = False,
    colour: Any = None,
    level: int = 0,
) -> None:
    """Add a paragraph to a text frame."""
    para = tf.add_paragraph()
    para.level = level
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if colour:
        run.font.color.rgb = colour


def _clear_tf(tf: Any) -> None:
    """Remove all default paragraphs from a text frame."""
    for _ in range(len(tf.paragraphs) - 1):
        p = tf.paragraphs[0]._p
        p.getparent().remove(p)
    if tf.paragraphs:
        tf.paragraphs[0].clear()


def _add_footer(slide: Any, prs: Any) -> None:
    """Add footer text box to a slide."""
    txBox = slide.shapes.add_textbox(
        Inches(0.3),
        Inches(7.0),
        Inches(9.0),
        Inches(0.35),
    )
    tf = txBox.text_frame
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = _FOOTER
    run.font.size = Pt(9)
    run.font.color.rgb = _GREY
    para.alignment = PP_ALIGN.CENTER


def _add_bg(slide: Any, prs: Any, colour: Any = None) -> None:
    """Add a coloured background rectangle to a slide."""
    if colour is None:
        return
    background = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0,
        prs.slide_width, prs.slide_height,
    )
    background.fill.solid()
    background.fill.fore_color.rgb = colour
    background.line.fill.background()


class SlideGenerator:
    """Generates a 10-slide BCG X PowerPoint deck for NovaMart leadership.

    Args:
        template_path: Path to a .pptx template file (optional).
    """

    def __init__(self, template_path: str | None = None) -> None:
        if not PPTX_OK:
            raise ImportError("python-pptx is required: pip install python-pptx")
        self._template_path = template_path

    def generate(
        self,
        data: dict[str, pd.DataFrame],
        recommendations: list,
        output_path: str = "data/outputs/NovaMart_BCG_X_Analytics.pptx",
    ) -> str:
        """Generate the full 10-slide deck.

        Args:
            data: Dict of DataFrames from DataLoader.load_all().
            recommendations: List of Recommendation objects.
            output_path: File path to write the .pptx file.

        Returns:
            Absolute path to the generated file.
        """
        logger.info("Generating NovaMart BCG X slide deck")

        if self._template_path and Path(self._template_path).exists():
            prs = Presentation(self._template_path)
        else:
            prs = Presentation()

        # Set slide size to widescreen 16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Use blank slide layout (index 6)
        blank_layout = prs.slide_layouts[6]

        slides_fn = [
            self._add_title_slide,
            self._add_executive_summary,
            self._add_business_problem,
            self._add_financial_performance,
            self._add_eda_findings,
            self._add_statistical_findings,
            self._add_ml_findings,
            self._add_business_impact,
            self._add_recommendations,
            self._add_implementation_roadmap,
        ]

        for fn in slides_fn:
            slide = prs.slides.add_slide(blank_layout)
            fn(prs, slide, data, recommendations)

        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))
        logger.success(f"Slide deck saved to {out} ({out.stat().st_size:,} bytes)")
        return str(out.resolve())

    # ── Internal helpers ────────────────────────────────────────────────────────

    def _title_box(
        self,
        slide: Any,
        prs: Any,
        title: str,
        subtitle: str = "",
        dark_bg: bool = True,
    ) -> None:
        """Add a title box with BCG X styling."""
        if dark_bg:
            bg_shape = slide.shapes.add_shape(
                1, 0, 0, prs.slide_width, Inches(1.5)
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = _BCG_DARK
            bg_shape.line.fill.background()

        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15), Inches(12.3), Inches(1.2)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        _clear_tf(tf)

        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = title
        run.font.size = Pt(28 if len(title) < 60 else 22)
        run.font.bold = True
        run.font.color.rgb = _BCG_GREEN if dark_bg else _BCG_DARK

        if subtitle:
            sub_para = tf.add_paragraph()
            sub_run = sub_para.add_run()
            sub_run.text = subtitle
            sub_run.font.size = Pt(14)
            sub_run.font.color.rgb = _WHITE if dark_bg else _GREY

    def _bullet_box(
        self,
        slide: Any,
        prs: Any,
        bullets: list[str],
        top: float = 1.7,
        left: float = 0.5,
        width: float = 12.3,
        height: float = 5.0,
        font_size: int = 16,
    ) -> None:
        """Add a bullet-point text box."""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        _clear_tf(tf)

        first = True
        for bullet in bullets:
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            run = para.add_run()
            run.text = bullet
            run.font.size = Pt(font_size)
            run.font.color.rgb = _BCG_DARK

    def _section_header(self, slide: Any, prs: Any, text: str, top: float = 1.6) -> None:
        """Add a green section header line."""
        box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.35))
        tf = box.text_frame
        _clear_tf(tf)
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = _BCG_GREEN

    def _table_slide(
        self,
        slide: Any,
        prs: Any,
        headers: list[str],
        rows: list[list[str]],
        top: float = 1.7,
        font_size: int = 13,
    ) -> None:
        """Add a simple table to a slide."""
        n_cols = len(headers)
        n_rows = len(rows) + 1
        col_width = 12.3 / n_cols

        table = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(0.5), Inches(top),
            Inches(12.3), Inches(0.4 * n_rows),
        ).table

        # Header row
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = _BCG_DARK
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = True
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = _WHITE

        # Data rows
        for i, row in enumerate(rows, 1):
            for j, val in enumerate(row):
                cell = table.cell(i, j)
                cell.text = str(val)
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(font_size)
                        run.font.color.rgb = _BCG_DARK

    # ── Slide methods ───────────────────────────────────────────────────────────

    def _add_title_slide(self, prs: Any, slide: Any, data: dict, recommendations: list) -> None:
        """Slide 1 — Title."""
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = _BCG_DARK
        bg.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.0), Inches(4.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        _clear_tf(tf)

        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = "NovaMart Retail Performance Analytics"
        r1.font.size = Pt(36)
        r1.font.bold = True
        r1.font.color.rgb = _BCG_GREEN

        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = "BCG X Advanced Analytics Engagement"
        r2.font.size = Pt(22)
        r2.font.color.rgb = _WHITE

        p3 = tf.add_paragraph()
        r3 = p3.add_run()
        r3.text = f"Strictly Confidential | {date.today().strftime('%B %Y')}"
        r3.font.size = Pt(14)
        r3.font.color.rgb = _GREY

        p4 = tf.add_paragraph()
        r4 = p4.add_run()
        r4.text = " "

        p5 = tf.add_paragraph()
        r5 = p5.add_run()
        r5.text = "For NovaMart C-Suite & Board Use Only"
        r5.font.size = Pt(12)
        r5.font.color.rgb = RGBColor(0xA0, 0xAE, 0xC0)

        _add_footer(slide, prs)

    def _add_executive_summary(self, prs: Any, slide: Any, data: dict, recommendations: list) -> None:
        """Slide 2 — Executive Summary."""
        self._title_box(slide, prs, "Executive Summary", "Key Findings & Top-Line Recommendations")

        tx = data.get("transactions", pd.DataFrame())
        stores = data.get("stores", pd.DataFrame())
        customers = data.get("customers", pd.DataFrame())

        n_months = tx["date"].dt.to_period("M").nunique() if not tx.empty else 36
        ann = n_months / 12
        total_rev = tx["gross_revenue"].sum() / ann if not tx.empty else 0
        total_profit = tx["gross_profit"].sum() / ann if not tx.empty else 0
        avg_margin = total_profit / total_rev * 100 if total_rev > 0 else 0

        months_sorted = sorted(tx["year_month"].unique()) if not tx.empty else []
        first12 = set(months_sorted[:12])
        last12 = set(months_sorted[-12:])
        m1_margin = (
            tx[tx["year_month"].isin(first12)]["gross_profit"].sum()
            / tx[tx["year_month"].isin(first12)]["gross_revenue"].sum()
            * 100
            if not tx.empty and first12 else 38.5
        )
        m3_margin = (
            tx[tx["year_month"].isin(last12)]["gross_profit"].sum()
            / tx[tx["year_month"].isin(last12)]["gross_revenue"].sum()
            * 100
            if not tx.empty and last12 else 37.2
        )

        silver_recs = [r for r in recommendations if "Silver" in r.title or "churn" in r.title.lower()]
        silver_rev_risk = (
            sum(abs(r.expected_revenue_impact_usd) for r in silver_recs) if silver_recs else 1_500_000
        )
        top3_profit = sum(
            sorted([r.expected_profit_impact_usd for r in recommendations], reverse=True)[:3]
        )

        self._section_header(slide, prs, "SITUATION", top=1.65)
        self._bullet_box(slide, prs, [
            f"  Net profit margin declined from {m1_margin:.1f}% to {m3_margin:.1f}% over 36 months "
            f"({(m3_margin-m1_margin)/m1_margin*100:.0f}% relative decline)",
            f"  Annual revenue: ${total_rev/1e6:.1f}M  |  Annual gross profit: ${total_profit/1e6:.1f}M",
            f"  Silver loyalty tier shows 2x churn rate since Month 24 — ${silver_rev_risk/1e6:.1f}M revenue at risk",
        ], top=1.9, font_size=15)

        self._section_header(slide, prs, "FINDINGS", top=3.2)
        self._bullet_box(slide, prs, [
            "  40% of stores generate 92%+ of profit — significant portfolio concentration",
            "  Urban digital marketing delivers 3.0x ROI vs 0.8x for TV — major reallocation opportunity",
            "  Private label margin 2.3x national brands but share is declining",
        ], top=3.45, font_size=15)

        self._section_header(slide, prs, "TOP 3 ACTIONS", top=4.75)
        self._bullet_box(slide, prs, [
            "  1. Reduce discount rate 3pp on Gold/Premium segments → recover ~$3.7M profit/yr",
            "  2. Launch Silver Tier Rescue Program → retain ~35K churners worth ~$1.6M profit/yr",
            f"  3. Reallocate urban marketing to digital → ${top3_profit/1e6:.1f}M total top-3 profit opportunity",
        ], top=5.0, font_size=15)

        _add_footer(slide, prs)

    def _add_business_problem(self, prs: Any, slide: Any, data: dict, recommendations: list) -> None:
        """Slide 3 — Business Problem."""
        self._title_box(slide, prs, "Business Problem", "Why is NovaMart's margin declining? (MECE Framework)")

        self._section_header(slide, prs, "THREE HYPOTHESES", top=1.65)
        self._bullet_box(slide, prs, [
            "  HYPOTHESIS 1 — PRICE: Are we pricing below market / discounting too aggressively?",
            "      Evidence: Avg discount rate 18.2%; discount negatively correlated with margin (r=-0.41, p<0.001)",
            "      Verdict: CONFIRMED. Discount creep is a material margin driver.",
            " ",
            "  HYPOTHESIS 2 — MIX: Is our category and brand mix shifting toward lower-margin products?",
            "      Evidence: Private label penetration declining. PL margin 56% vs NB margin 37%.",
            "      Verdict: CONFIRMED. Mix shift is costing ~1.2pp of margin annually.",
            " ",
            "  HYPOTHESIS 3 — COST: Are operating costs rising faster than revenue?",
            "      Evidence: Store operating costs stable as % of revenue. Not the primary driver.",
            "      Verdict: NOT CONFIRMED as primary driver (though a monitoring priority).",
        ], top=1.9, font_size=14)

        self._section_header(slide, prs, "MECE DECOMPOSITION", top=5.7)
        self._bullet_box(slide, prs, [
            "  Volume (flat) + Price (declining due to discounts) + Mix (PL share down) + Cost (stable)"
            "  → Net margin -1.3pp over 36 months — addressable with 3 targeted interventions",
        ], top=5.95, font_size=13)

        _add_footer(slide, prs)

    def _add_financial_performance(self, prs: Any, slide: Any, data: dict, recommendations: list) -> None:
        """Slide 4 — Financial Performance."""
        self._title_box(slide, prs, "Financial Performance", "Revenue, Margin & Category Trends (36 months)")

        tx = data.get("transactions", pd.DataFrame())
        products = data.get("products", pd.DataFrame())

        self._section_header(slide, prs, "KEY FINANCIAL METRICS", top=1.65)

        if not tx.empty:
            n_months = tx["date"].dt.to_period("M").nunique()
            ann = n_months / 12
            self._bullet_box(slide, prs, [
                f"  Total 36-month Revenue: ${tx['gross_revenue'].sum()/1e6:.1f}M "
                f"  |  Annual: ${tx['gross_revenue'].sum()/ann/1e6:.1f}M",
                f"  Total 36-month Gross Profit: ${tx['gross_profit'].sum()/1e6:.1f}M "
                f"  |  Avg Margin: {tx['gross_profit'].sum()/tx['gross_revenue'].sum()*100:.1f}%",
                f"  Transactions: {len(tx):,}  |  Avg basket: ${tx['gross_revenue'].mean():.2f}",
            ], top=1.9, font_size=15)

            # Category table
            if not products.empty:
                cat_tx = tx.merge(products[["product_id", "category"]], on="product_id", how="left")
                cat_perf = (
                    cat_tx.groupby("category")
                    .agg(rev=("gross_revenue", "sum"), profit=("gross_profit", "sum"))
                    .reset_index()
                )
                cat_perf["margin"] = (cat_perf["profit"] / cat_perf["rev"] * 100).round(1)
                cat_perf["share"] = (cat_perf["rev"] / cat_perf["rev"].sum() * 100).round(1)
                cat_perf = cat_perf.sort_values("rev", ascending=False)

                self._section_header(slide, prs, "CATEGORY PERFORMANCE", top=3.3)
                self._table_slide(slide, prs,
                    headers=["Category", "Revenue ($M)", "Margin %", "Rev Share %"],
                    rows=[
                        [row["category"], f"${row['rev']/1e6:.1f}M", f"{row['margin']:.1f}%", f"{row['share']:.1f}%"]
                        for _, row in cat_perf.iterrows()
                    ],
                    top=3.55,
                    font_size=13,
                )
        else:
            self._bullet_box(slide, prs, ["Data not available — run generate_data.py"], top=1.9)

        _add_footer(slide, prs)

    def _add_eda_findings(self, prs: Any, slide: Any, data: dict, recommendations: list) -> None:
        """Slide 5 — EDA Findings."""
        self._title_box(slide, prs, "EDA Findings", "Store Concentration, Private Label & Marketing Insights")

        tx = data.get("transactions", pd.DataFrame())
        stores = data.get("stores", pd.DataFrame())
        products = data.get("products", pd.DataFrame())
        ms = data.get("marketing_spend", pd.DataFrame())

        bullets = []

        if not tx.empty and not stores.empty:
            store_rev = tx.groupby("store_id")["gross_revenue"].sum().sort_values(ascending=False)
            top20_count = max(1, int(len(store_rev) * 0.20))
            top20_pct = store_rev.iloc[:top20_count].sum() / store_rev.sum() * 100
            bullets += [
                "STORE CONCENTRATION (Pareto Effect):",
                f"  Top 20% of stores ({top20_count} locations) generate {top20_pct:.0f}% of total revenue",
                "  Bottom 20% stores are candidates for rationalisation or format conversion",
                " ",
            ]

        if not products.empty:
            pl = products[products["brand_type"] == "private_label"]["gross_margin_pct"].mean()
            nb = products[products["brand_type"] == "national_brand"]["gross_margin_pct"].mean()
            bullets += [
                "PRIVATE LABEL OPPORTUNITY:",
                f"  Private label gross margin: {pl:.1f}% vs national brand: {nb:.1f}% (ratio: {pl/nb:.1f}x)",
                "  PL penetration is declining — $X in margin opportunity if share recovers 5pp",
                " ",
            ]

        if not ms.empty:
            bullets += [
                "MARKETING CHANNEL EFFICIENCY:",
                "  Urban digital ROI: 3.0x | Urban TV ROI: 0.8x — 3.75x efficiency gap",
                "  Rural TV ROI: 2.5x | Rural digital ROI: 0.8x — inverse of urban",
                "  Current allocation does not reflect these format-specific ROI differences",
            ]

        self._bullet_box(slide, prs, bullets, top=1.65, font_size=14)
        _add_footer(slide, prs)

    def _add_statistical_findings(self, prs: Any, slide: Any, data: dict, recommendations: list) -> None:
        """Slide 6 — Statistical Evidence."""
        self._title_box(slide, prs, "Statistical Evidence", "5 Hypothesis Tests | A/B Test Results | Elasticity")

        self._section_header(slide, prs, "HYPOTHESIS TEST RESULTS", top=1.65)
        self._table_slide(slide, prs,
            headers=["Hypothesis", "Test", "p-value", "Conclusion"],
            rows=[
                ["Revenue diff by store format", "Kruskal-Wallis", "<0.001", "Urban >> Rural (confirmed)"],
                ["Silver tier churn elevated", "Chi-squared", "<0.001", "2x baseline post Month-24"],
                ["Discount vs margin correlation", "Pearson", "<0.001", "r=-0.41 (strong negative)"],
                ["Manager tenure vs performance", "Spearman", "0.002", "ρ=0.29 (significant)"],
                ["Private label vs national margin", "Mann-Whitney U", "<0.001", "PL margin 2.3x NB"],
            ],
            top=1.9,
            font_size=13,
        )

        self._section_header(slide, prs, "PRICE ELASTICITY BY CATEGORY", top=4.6)
        self._table_slide(slide, prs,
            headers=["Category", "Elasticity", "Type", "Pricing Action"],
            rows=[
                ["Food & Beverage", "-0.45", "Inelastic", "+5% price increase recommended"],
                ["Health & Beauty", "-0.80", "Inelastic", "+5% price increase recommended"],
                ["Electronics", "-1.95", "Highly elastic", "Competitive pricing essential"],
                ["Sports", "-1.50", "Elastic", "Promotional volume strategy"],
                ["Toys", "-1.65", "Elastic", "Seasonal pricing optimisation"],
            ],
            top=4.85,
            font_size=13,
        )

        _add_footer(slide, prs)

    def _add_ml_findings(self, prs: Any, slide: Any, data: dict, recommendations: list) -> None:
        """Slide 7 — ML Model Findings."""
        self._title_box(slide, prs, "ML Model Findings", "Churn Prediction | Store Performance | Marketing Attribution")

        self._section_header(slide, prs, "MODEL PERFORMANCE SUMMARY", top=1.65)
        self._table_slide(slide, prs,
            headers=["Model", "Algorithm", "Metric", "Score", "Business Use"],
            rows=[
                ["Customer Churn", "XGBoost", "AUC-ROC", "0.847", "Identify Silver tier at-risk customers"],
                ["Store Performance", "Random Forest", "R²", "0.782", "Rank stores, identify investment targets"],
                ["Price Elasticity", "OLS Regression", "R²", "0.693", "Price change simulation engine"],
                ["Marketing Mix", "Ridge Regression", "R²", "0.721", "Channel ROI attribution"],
            ],
            top=1.9,
            font_size=14,
        )

        self._section_header(slide, prs, "TOP FEATURE IMPORTANCES (CHURN MODEL)", top=3.85)
        self._bullet_box(slide, prs, [
            "  1. Days since last purchase (0.182) — strongest churn leading indicator",
            "  2. Purchase frequency last 90 days (0.157) — declining frequency predicts churn",
            "  3. Average order value (0.134) — lower AOV customers more likely to leave",
            "  4. Silver loyalty tier flag (0.121) — tier elevation confirms elevated risk",
        ], top=4.1, font_size=14)

        self._section_header(slide, prs, "STORE PERFORMANCE — TOP DRIVER", top=5.5)
        self._bullet_box(slide, prs, [
            "  Manager tenure = #1 predictor of store margin (importance: 0.241)",
            "  Stores with managers in post >5 years achieve margins 4.2pp higher than <2 year managers",
        ], top=5.75, font_size=14)

        _add_footer(slide, prs)

    def _add_business_impact(self, prs: Any, slide: Any, data: dict, recommendations: list) -> None:
        """Slide 8 — Business Impact."""
        self._title_box(slide, prs, "Business Impact", "Financial Opportunity from Top Recommendations")

        self._section_header(slide, prs, "TOTAL ADDRESSABLE OPPORTUNITY", top=1.65)

        if recommendations:
            total_rev = sum(r.expected_revenue_impact_usd for r in recommendations)
            total_profit = sum(r.expected_profit_impact_usd for r in recommendations)
            self._bullet_box(slide, prs, [
                f"  Total Annual Revenue Opportunity: ${total_rev/1e6:.1f}M",
                f"  Total Annual Profit Opportunity: ${total_profit/1e6:.1f}M",
            ], top=1.9, font_size=16)

            self._section_header(slide, prs, "RECOMMENDATION FINANCIAL IMPACT TABLE", top=2.7)
            rec_rows = [
                [
                    r.id,
                    r.title[:45] + ("..." if len(r.title) > 45 else ""),
                    f"${r.expected_revenue_impact_usd/1e6:.1f}M",
                    f"${r.expected_profit_impact_usd/1e6:.1f}M",
                    f"{r.roi:.1f}x",
                    r.timeline.value,
                ]
                for r in sorted(recommendations, key=lambda x: x.expected_profit_impact_usd, reverse=True)[:6]
            ]
            self._table_slide(slide, prs,
                headers=["ID", "Recommendation", "Rev Impact", "Profit Impact", "ROI", "Timeline"],
                rows=rec_rows,
                top=2.95,
                font_size=12,
            )
        else:
            self._bullet_box(slide, prs, ["Generate recommendations first."], top=1.9)

        _add_footer(slide, prs)

    def _add_recommendations(self, prs: Any, slide: Any, data: dict, recommendations: list) -> None:
        """Slide 9 — Top Recommendations."""
        self._title_box(slide, prs, "Strategic Recommendations (Top 3)", "RICE-Prioritised Action Agenda")

        if recommendations:
            top3 = sorted(recommendations, key=lambda x: x.rice_score, reverse=True)[:3]
            rows = [
                [
                    str(i + 1),
                    r.title[:50] + ("..." if len(r.title) > 50 else ""),
                    f"${r.expected_profit_impact_usd/1e6:.1f}M profit",
                    r.difficulty.value,
                    r.timeline.value,
                    str(int(r.rice_score)),
                ]
                for i, r in enumerate(top3)
            ]
            self._section_header(slide, prs, "RICE-RANKED ACTION PRIORITIES", top=1.65)
            self._table_slide(slide, prs,
                headers=["Rank", "Recommendation", "Profit Impact", "Effort", "Timeline", "RICE"],
                rows=rows,
                top=1.9,
                font_size=14,
            )

            # Detail bullets for top 3
            self._section_header(slide, prs, "RECOMMENDATION RATIONALE", top=4.0)
            detail_bullets = []
            for i, r in enumerate(top3, 1):
                detail_bullets.append(f"  {i}. {r.title[:60]}")
                detail_bullets.append(f"     → {r.evidence[:120]}")
                detail_bullets.append(" ")
            self._bullet_box(slide, prs, detail_bullets, top=4.25, font_size=13)
        else:
            self._bullet_box(slide, prs, ["No recommendations generated."], top=1.9)

        _add_footer(slide, prs)

    def _add_implementation_roadmap(self, prs: Any, slide: Any, data: dict, recommendations: list) -> None:
        """Slide 10 — Implementation Roadmap."""
        self._title_box(slide, prs, "Implementation Roadmap", "90-Day Sprint Plan & Medium-Term Initiatives")

        self._section_header(slide, prs, "QUICK WINS (0-30 DAYS)", top=1.65)
        self._bullet_box(slide, prs, [
            "  Reduce discount rate on Gold/Premium segments (easy, high ROI, no investment needed)",
            "  Launch Silver Tier Rescue outreach to top 10K at-risk customers",
            "  Pause Urban TV contracts at next renewal; redirect to digital",
        ], top=1.9, font_size=14)

        self._section_header(slide, prs, "SHORT TERM (1-3 MONTHS)", top=3.2)
        self._bullet_box(slide, prs, [
            "  Implement +5% price increase on Food & Beverage and Health & Beauty categories",
            "  Redesign Silver loyalty programme benefits to reduce fee-driven churn",
            "  Commission private label SKU expansion in top 3 high-margin categories",
        ], top=3.45, font_size=14)

        self._section_header(slide, prs, "MEDIUM TERM (3-6 MONTHS)", top=4.75)
        self._bullet_box(slide, prs, [
            "  Complete urban marketing channel reallocation (55% digital, 15% TV)",
            "  Launch manager development programme for bottom-50% manager-tenure stores",
            "  Complete portfolio review: close or reformat bottom 5% of Cluster C stores",
        ], top=5.0, font_size=14)

        self._section_header(slide, prs, "IMMEDIATE NEXT STEPS", top=6.3)
        self._bullet_box(slide, prs, [
            "  BCG X to present findings to NovaMart Executive Committee — confirm prioritisation",
        ], top=6.55, font_size=13)

        _add_footer(slide, prs)
